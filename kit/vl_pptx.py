"""vl_pptx — le pont entre la bibliothèque visual-lab et un vrai .pptx.

POURQUOI CE FICHIER EXISTE
Un pattern de visual-lab est un fragment HTML : on ne le colle pas dans un .pptx. Sans pont,
la bibliothèque ne servirait qu'au web et `deck-builder` continuerait à réinventer chaque
vignette. Le pont ne duplique PAS le pattern : il lit `index.json` (généré par
`node bin/index.mjs`) et n'en reprend que ce qui voyage — les RATIOS et les tokens. Une seule
source de vérité, deux rendus.

LES PIXELS NE VOYAGENT PAS, LES RAPPORTS OUI
Les fragments sont réglés pour une slide de 1600 px de large. Converti bêtement, un corps de
17 px donne 10,2 pt : illisible en projection, et sous le plancher de `deck-builder`
(FLOOR_BODY_PT = 14). On ancre donc l'échelle sur le CORPS en points (14 par défaut) et on
dérive tout le reste des rapports du pattern. La vignette garde sa tension typographique et
reste lisible — c'est le seul réglage qui satisfait les deux systèmes.

USAGE
    import sys; sys.path.insert(0, str(Path.home() / "visual-lab" / "kit"))
    from vl_pptx import stat_block, brief_card, audit

    m = stat_block(slide, 1.0, 2.2, 3.6, 2.8, "POTENTIAL", "156%",
                   "Compound growth of the addressable segment.")
    audit([m])          # contrôle MATHÉMATIQUE ; lève si un ratio du pattern n'est pas tenu

Contrôle VISUEL ensuite, jamais à la place : `deck-builder/render_check.py`.
"""

from __future__ import annotations

import json
import math
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

LIB = Path.home() / "visual-lab"
_INDEX_CACHE: dict | None = None

# Plancher de lisibilité : identique à deck-builder.FLOOR_BODY_PT. Dupliqué comme CONSTANTE et
# non importé, pour que ce kit reste utilisable hors de ce skill — mais les deux doivent bouger
# ensemble.
FLOOR_BODY_PT = 14.0


# ————————————————————————— lecture de la bibliothèque —————————————————————————

def index() -> dict:
    """Charge `index.json`. Il est REGÉNÉRÉ par `node bin/index.mjs` : s'il manque, c'est
    l'indexation qui n'a pas été relancée, pas le kit qui est cassé."""
    global _INDEX_CACHE
    if _INDEX_CACHE is None:
        path = LIB / "index.json"
        if not path.exists():
            raise FileNotFoundError(
                f"{path} absent — lance `cd {LIB} && node bin/index.mjs` "
                "(l'index est généré, il n'est pas versionné à la main)."
            )
        _INDEX_CACHE = json.loads(path.read_text())
    return _INDEX_CACHE


def pattern(pid: str) -> dict:
    for p in index()["patterns"]:
        if p["id"] == pid:
            return p
    known = ", ".join(p["id"] for p in index()["patterns"])
    raise KeyError(f"pattern inconnu : {pid}\nDisponibles : {known}")


def geom(pid: str, key: str):
    """Un champ de `geometry`, avec l'erreur qui dit quoi faire. Le cas courant n'est pas un
    pattern mal écrit mais un `index.json` PÉRIMÉ : le .json a été enrichi et l'indexation
    n'a pas été relancée. Le kit doit le dire, pas sortir un KeyError nu."""
    g = pattern(pid).get("geometry") or {}
    if key not in g:
        raise KeyError(
            f"{pid} : geometry.{key} absent de index.json. "
            f"Si patterns/{pid}.json le porte, l'index est périmé → "
            f"`cd {LIB} && node bin/index.mjs`."
        )
    return g[key]


def tokens(pid: str) -> dict:
    """Les tokens de la référence du pattern, en hexadécimal."""
    pat = pattern(pid)
    for s in index()["refs"]:
        if s["id"] == pat["ref"]:
            return s["tokens"]
    return {}


def rgb(pid: str, token: str, fallback: str = "#000000") -> RGBColor:
    value = tokens(pid).get(token, fallback)
    return RGBColor.from_string(value.lstrip("#")[:6].upper())


def scale(pid: str, body_pt: float = FLOOR_BODY_PT) -> dict[str, float]:
    """Corps en POINTS pour chaque slot, dérivés des rapports du fragment et ancrés sur le
    corps. Le plancher s'applique au corps ET aux libellés : un micro-libellé qui tombe à 8 pt
    n'est plus une micro-capitale, c'est une note de bas de page illisible."""
    px = geom(pid, "type_px")
    # L'ancre est le CORPS, sous le nom qu'il porte dans ce pattern-là. `note` en fait partie :
    # sans elle, un pattern dont le corps s'appelle ainsi s'ancrait sur son plus GROS slot, et
    # tout le reste tombait sous le plancher pour y être remonté — deux niveaux de typo écrasés
    # en un seul, la hiérarchie perdue à la conversion.
    for name in ("body", "desc", "note"):
        if name in px:
            anchor_px, anchor_name = px[name], name
            break
    else:  # un pattern de titre : le titre EST l'ancre
        anchor_name = next(iter(px))
        anchor_px = px[anchor_name]
    factor = body_pt / anchor_px
    out = {k: round(v * factor, 1) for k, v in px.items()}
    for k in out:
        if k not in ("kicker", "label") and out[k] < body_pt:
            out[k] = body_pt
    return out


def contrast(fg: str, bg: str) -> float:
    """Rapport de contraste WCAG entre deux couleurs hexadécimales."""
    def lum(h: str) -> float:
        h = h.lstrip("#")
        parts = [int(h[i:i + 2], 16) / 255 for i in (0, 2, 4)]
        parts = [c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4 for c in parts]
        return 0.2126 * parts[0] + 0.7152 * parts[1] + 0.0722 * parts[2]

    a, b = lum(fg), lum(bg)
    hi, lo = max(a, b), min(a, b)
    return (hi + 0.05) / (lo + 0.05)


# ————————————————————————— la primitive : le coin chanfreiné —————————————————————————

def destyle(shape):
    """Coupe le lien au thème : aucune ombre, aucun effet hérité.

    ⚠️ Piège vérifié au rendu LibreOffice : `shape.shadow.inherit = False` écrit bien un
    `<a:effectLst/>` vide, et ça ne suffit PAS. python-pptx ajoute aussi un `<p:style>` qui
    référence les effets du thème (`effectRef`) — le moteur de rendu l'applique et la carte
    sort avec une ombre portée grise en bas à droite. Sur une charte à angles vifs, cette
    ombre est un défaut visible (constaté sur le premier .pptx de démonstration). Il faut
    RETIRER l'élément `<p:style>`, pas seulement neutraliser l'effet.

    ⚠️ Et ça vaut pour les TRAITS autant que pour les aplats : un connecteur naît lui aussi
    avec son `<p:style>`, et un filet de rappel de 0,75 pt doublé d'une ombre grise se lit
    comme un trait sale (constaté sur le premier rendu de `layer_stack`).
    """
    shape.shadow.inherit = False
    for style in shape._element.findall(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style"
    ):
        shape._element.remove(style)
    return shape


def flatten(shape, fill: RGBColor):
    """Aplat NET : couleur pleine, aucun contour, aucune ombre."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return destyle(shape)


def notched_rect(slide, x, y, w, h, fill: RGBColor, corner: str = "br",
                 notch_ratio: float = 0.075):
    """Rectangle à UN coin coupé à 45°, en freeform 5 points.

    ⚠️ Ne pas chercher à obtenir ce coin avec `MSO_SHAPE.ROUND_1_RECTANGLE` ni un
    `SNIP_1_RECTANGLE` : le snip de PowerPoint est piloté par un ajustement en pourcentage de
    la PLUS PETITE dimension, donc la coupe change de longueur dès que la carte change de
    proportion — et la charte devient méconnaissable d'une slide à l'autre. Le freeform fixe la
    coupe en absolu, comme le `clip-path` du fragment.
    """
    x, y, w, h = Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h))
    n = int(w * notch_ratio)
    corners = {
        "br": [(0, 0), (w, 0), (w, h - n), (w - n, h), (0, h)],
        "tr": [(0, 0), (w - n, 0), (w, n), (w, h), (0, h)],
        "bl": [(0, 0), (w, 0), (w, h), (n, h), (0, h - n)],
        "tl": [(n, 0), (w, 0), (w, h), (0, h), (0, n)],
    }
    if corner not in corners:
        raise ValueError(f"corner doit être br/tr/bl/tl, reçu {corner!r}")
    pts = [(x + dx, y + dy) for dx, dy in corners[corner]]
    builder = slide.shapes.build_freeform(pts[0][0], pts[0][1])
    builder.add_line_segments(pts[1:], close=True)
    return flatten(builder.convert_to_shape(), fill)


def _text(slide, x, y, w, h, runs, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """Boîte de texte à marges NULLES. La marge interne par défaut de python-pptx (0,05\" de
    chaque côté) suffit à décoller un filet censé toucher la lettre, ou à faire déborder un
    corps calé au pixel — on l'annule toujours."""
    box = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (text, size_pt, color, bold, spacing) in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
        if spacing:
            # l'interlettrage n'est pas exposé par python-pptx : on l'écrit dans le XML
            run.font._rPr.set("spc", str(int(spacing * size_pt * 100)))
    return box


# ————————————————————————— les vignettes —————————————————————————

def stat_block(slide, x_in, y_in, w_in, h_in, kicker: str, figure: str, body: str,
               pid: str = "card-03-stat-accent", body_pt: float = FLOOR_BODY_PT) -> dict:
    """Vignette statistique sur aplat d'accent (cf. `card-03-stat-accent`).
    Retourne les mesures pour `audit()`."""
    x, y, w, h = (Inches(v) for v in (x_in, y_in, w_in, h_in))
    pads = geom(pid, "pad_ratio")
    pad_x, pad_t, pad_b = int(w * pads["x"]), int(w * pads["top"]), int(w * pads["bottom"])
    notch_ratio = geom(pid, "ratios")["notch_over_width"]
    sizes = scale(pid, body_pt)

    fill = rgb(pid, "--vl-orange")
    ink = rgb(pid, "--vl-white", "#FFFFFF")
    ink_body = rgb(pid, "--vl-ink-muted-invert", "#FFE0D3")
    notched_rect(slide, x, y, w, h, fill, corner="br", notch_ratio=notch_ratio)

    inner_w = w - 2 * pad_x
    _text(slide, x + pad_x, y + pad_t, inner_w, Pt(sizes["kicker"] * 1.4),
          [(kicker.upper(), sizes["kicker"], ink, True, 0.14)])
    _text(slide, x + pad_x, y + pad_t + Pt(sizes["kicker"] * 2.2), inner_w,
          Pt(sizes["figure"] * 1.1), [(figure, sizes["figure"], ink, True, 0)])
    # Le corps est calé EN BAS : le vide entre le chiffre et lui sépare la donnée de son
    # commentaire (cf. les notes du pattern) — il n'est pas à combler.
    body_h = Pt(sizes["body"] * 1.55 * 3)
    body_box = _text(slide, x + pad_x, y + h - pad_b - body_h, int(inner_w * 0.82), body_h,
                     [(body, sizes["body"], ink_body, False, 0)], anchor=MSO_ANCHOR.BOTTOM)

    return {
        "pid": pid, "w": w, "h": h,
        "measured": {
            "notch_over_width": notch_ratio,
            "pad_x_over_width": pad_x / w,
            "figure_over_kicker": sizes["figure"] / sizes["kicker"],
            "body_top_over_height": (body_box.top - y) / h,
        },
        "type_pt": sizes,
        "contrasts": {
            "figure": contrast(tokens(pid).get("--vl-white", "#FFFFFF"), tokens(pid)["--vl-orange"]),
            "body": contrast(tokens(pid).get("--vl-ink-muted-invert", "#FFE0D3"),
                             tokens(pid)["--vl-orange"]),
        },
        "min_contrast": {"figure": 3.0, "body": 3.0},
    }


def brief_card(slide, x_in, y_in, w_in, h_in, label: str, body: str,
               pid: str = "card-02-notched-brief", body_pt: float = FLOOR_BODY_PT) -> dict:
    """Carte chanfreinée de second rang (cf. `card-02-notched-brief`). Le chanfrein est en
    HAUT-droit : c'est l'orientation, pas la couleur, qui la distingue de la vignette
    d'accent — un tirage N&B doit encore faire la différence."""
    x, y, w, h = (Inches(v) for v in (x_in, y_in, w_in, h_in))
    pads = geom(pid, "pad_ratio")
    pad_x, pad_t, pad_b = int(w * pads["x"]), int(w * pads["top"]), int(w * pads["bottom"])
    notch_ratio = geom(pid, "ratios")["notch_over_width"]
    sizes = scale(pid, body_pt)

    notched_rect(slide, x, y, w, h, rgb(pid, "--vl-steel"), corner="tr",
                 notch_ratio=notch_ratio)
    inner_w = w - 2 * pad_x
    _text(slide, x + pad_x, y + pad_t, int(inner_w * 0.7), Pt(sizes["label"] * 2.9),
          [(label.upper(), sizes["label"], rgb(pid, "--vl-black"), True, 0.14)])
    body_h = Pt(sizes["body"] * 1.55 * 3)
    body_box = _text(slide, x + pad_x, y + h - pad_b - body_h, inner_w, body_h,
                     [(body, sizes["body"], rgb(pid, "--vl-ink-on-steel", "#3E4746"), False, 0)],
                     anchor=MSO_ANCHOR.BOTTOM)

    return {
        "pid": pid, "w": w, "h": h,
        "measured": {
            "notch_over_width": notch_ratio,
            "pad_x_over_width": pad_x / w,
            "body_top_over_height": (body_box.top - y) / h,
        },
        "type_pt": sizes,
        "contrasts": {
            "body": contrast(tokens(pid).get("--vl-ink-on-steel", "#3E4746"),
                             tokens(pid)["--vl-steel"]),
            "label": contrast(tokens(pid)["--vl-black"], tokens(pid)["--vl-steel"]),
        },
        "min_contrast": {"body": 4.5, "label": 7.0},
    }


def title_leading_rule(slide, x_in, y_in, w_in, lines: list[str],
                       pid: str = "title-leading-rule", title_pt: float = 34.0) -> dict:
    """Titre en capitales précédé du filet d'accent COLLÉ (cf. `title-leading-rule`).
    Le filet est un rectangle plein de la hauteur du bloc de texte, posé à décalage nul."""
    rule_w = int(Pt(title_pt) * geom(pid, "ratios")["rule_width_over_font_size"])
    x, y, w = Inches(x_in), Inches(y_in), Inches(w_in)
    block_h = int(Pt(title_pt * 1.05) * len(lines))

    rule = flatten(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
                               Emu(rule_w), Emu(block_h)),
        rgb(pid, "--vl-orange"),
    )

    box = _text(slide, x + rule_w, y, w - rule_w, block_h,
                [(l.upper(), title_pt, rgb(pid, "--vl-black"), True, 0) for l in lines])
    for para in box.text_frame.paragraphs:
        para.line_spacing = 1.05

    return {
        "pid": pid, "w": w, "h": block_h,
        "measured": {"rule_width_over_font_size": rule_w / Pt(title_pt),
                     "rule_gap_pt": (box.left - (rule.left + rule.width)) / 12700},
        "type_pt": {"title": title_pt},
        "contrasts": {"title": contrast(tokens(pid)["--vl-black"], tokens(pid)["--vl-bg"])},
        "min_contrast": {"title": 7.0},
    }


def _alpha(shape, value: float):
    """Opacité d'un aplat déjà rempli.

    ⚠️ python-pptx n'expose AUCUNE transparence : ni `fill.transparency`, ni un argument de
    `fore_color`. Elle s'écrit dans le XML du remplissage — `<a:alpha val="44000"/>` pour 44 %.
    Sans elle, une pile de couches sort en trois aplats identiques et le rang, qui EST
    l'information du schéma, disparaît.
    """
    if value >= 1:
        return shape
    fill = shape._element.spPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    )
    srgb = fill.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    srgb.append(srgb.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
        {"val": str(int(round(value * 100000)))},
    ))
    return shape


def layer_stack(slide, x_in, y_in, w_in, layers: list[tuple[str, str]],
                pid: str = "diagram-01-layer-stack", body_pt: float = FLOOR_BODY_PT) -> dict:
    """Pile de couches isométriques légendées (cf. `diagram-01-layer-stack`).

    `layers` va du HAUT vers le BAS — l'ordre est l'information. La HAUTEUR n'est pas un
    argument : elle est dictée par la géométrie (padding + hauteur de plan + (n−1) × pas), donc
    la donner reviendrait à pouvoir écraser la pile. Elle est retournée dans la mesure.

    Le losange est un `MSO_SHAPE.DIAMOND` et non un freeform : contrairement au chanfrein de
    `notched_rect`, le diamant de PowerPoint EST exactement le `clip-path` du fragment
    (polygon 50/0 100/50 50/100 0/50), il n'a pas d'ajustement qui dérive avec la proportion.
    """
    if not 3 <= len(layers) <= 5:
        raise ValueError(
            f"{pid} attend 3 à 5 couches, reçu {len(layers)} — en dessous la pile ne dit pas "
            "un ordre, au-dessus les plans du bas s'effacent avant d'être comptés "
            "(cf. `avoid_when` du pattern)."
        )
    n = len(layers)
    x, y, w = Inches(x_in), Inches(y_in), Inches(w_in)
    r, pads = geom(pid, "ratios"), geom(pid, "pad_ratio")
    sizes = scale(pid, body_pt)

    plane_w = int(w * r["plane_w_over_width"])
    plane_h = int(plane_w * r["plane_h_over_w"])
    pitch = int(plane_h * r["pitch_over_plane_h"])
    pad_t, pad_b = int(w * pads["top"]), int(w * pads["bottom"])
    height = pad_t + plane_h + (n - 1) * pitch + pad_b

    # Le PAS est géométrique (il vient des ratios) ; le CORPS, lui, est ancré sur le plancher de
    # lisibilité en points. Les deux ne tiennent ensemble qu'au-delà d'une certaine largeur : en
    # dessous, l'explication d'une couche mord sur le nom de la suivante — et rien ne le
    # signalerait, un texte qui déborde n'étant pas une erreur en .pptx. Le budget de texte est
    # celui du fragment : la ligne de titre, son interligne, puis DEUX lignes d'explication.
    name_h = Pt(sizes["name"] * 1.2)
    note_gap = Pt(sizes["note"] * 0.6)
    text_h = name_h + note_gap + 2 * Pt(sizes["note"] * 1.7)
    if pitch < text_h:
        need = text_h / (r["pitch_over_plane_h"] * r["plane_h_over_w"] * r["plane_w_over_width"])
        # arrondi au DIXIÈME SUPÉRIEUR : une largeur conseillée qu'on ne peut pas recopier telle
        # quelle sans relever la même erreur ne conseille rien.
        need_in = math.ceil(need / 914400 * 10) / 10
        raise ValueError(
            f"{pid} : à {w_in:.2f}\" de large, le pas ({pitch / 12700:.1f} pt) est plus court "
            f"que le texte d'une couche ({text_h / 12700:.1f} pt à {body_pt} pt de corps) — les "
            f"rangées se chevaucheraient. Élargir à {need_in}\" au moins, ou baisser "
            "`body_pt` (au risque du plancher de lisibilité)."
        )

    rule_x = x + int(plane_w * r["rule_start_over_plane_w"])
    label_x = x + int(w * r["label_x_over_width"])
    rule_end = label_x - int(w * r["rule_gap_over_width"])

    # L'opacité suit une progression GÉOMÉTRIQUE. Les trois premiers alphas sont des tokens du
    # système ; au-delà on POURSUIT la raison relevée entre les deux premiers, on ne l'invente
    # pas — un pas soustractif rendrait la pile plate (cf. les notes du pattern).
    tok = tokens(pid)
    alphas = [float(tok.get(f"--vl-layer-alpha-{i + 1}", 0)) for i in range(3)]
    ratio = alphas[1] / alphas[0]
    while len(alphas) < n:
        alphas.append(alphas[-1] * ratio)
    alphas = alphas[:n]

    coral, ink = rgb(pid, "--vl-coral"), rgb(pid, "--vl-ink")
    ink_muted, hairline = rgb(pid, "--vl-ink-muted"), rgb(pid, "--vl-hairline")

    # L'ordre d'insertion EST l'ordre d'empilement. Il porte deux choses du fragment, et pas
    # seulement l'esthétique : les filets d'ABORD, parce qu'ils naissent SOUS les losanges et
    # que leur longueur visible est dictée par la forme (posés au-dessus, ils barreraient le
    # plan de tête) ; puis les plans du BAS vers le haut, la couche de tête devant.
    centers = [y + pad_t + i * pitch + plane_h / 2 for i in range(n)]
    for center in centers:
        rule = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(int(rule_x)), Emu(int(center)),
            Emu(int(rule_end)), Emu(int(center)))
        rule.line.color.rgb = hairline
        rule.line.width = Pt(0.75)
        destyle(rule)

    for i in range(n - 1, -1, -1):
        plane = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, Emu(int(x)), Emu(int(y + pad_t + i * pitch)),
            Emu(plane_w), Emu(plane_h))
        _alpha(flatten(plane, coral), alphas[i])

    for i, (name, note) in enumerate(layers):
        center = centers[i]
        _text(slide, label_x, center - name_h / 2, w - (label_x - x), name_h,
              [(name.upper(), sizes["name"], ink, False, 0.13)], anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, label_x, center + name_h / 2 + note_gap,
              w - (label_x - x), text_h - name_h - note_gap,
              [(note.upper(), sizes["note"], ink_muted, False, 0.05)])

    return {
        "pid": pid, "w": w, "h": height,
        "measured": {
            "plane_w_over_width": plane_w / w,
            "plane_h_over_w": plane_h / plane_w,
            "pitch_over_plane_h": pitch / plane_h,
            "rule_start_over_plane_w": (rule_x - x) / plane_w,
            "label_x_over_width": (label_x - x) / w,
            "rule_gap_over_width": (label_x - rule_end) / w,
        },
        "type_pt": sizes,
        "contrasts": {
            "name": contrast(tok["--vl-ink"], tok["--vl-bg"]),
            "note": contrast(tok["--vl-ink-muted"], tok["--vl-bg"]),
        },
        "min_contrast": {"name": 7.0, "note": 4.5},
    }


# ————————————————————————— le contrôle mathématique —————————————————————————

def audit(measures: list[dict], tol: float = 0.02, raise_on_fail: bool = True) -> list[dict]:
    """Confronte ce qui a été RÉELLEMENT posé sur la slide aux ratios du pattern, aux
    contrastes minimaux et au plancher de lisibilité. C'est le pendant .pptx de
    `node bin/check.mjs` : mêmes ratios, même exigence.

    Retourne la liste des constats ; lève `AssertionError` si l'un échoue (le geste attendu est
    de corriger et de relancer, pas de livrer avec un écart connu)."""
    out = []
    for m in measures:
        want = pattern(m["pid"])["geometry"]["ratios"]
        for key, got in m["measured"].items():
            if key in want:
                ok = abs(got - want[key]) <= max(tol, abs(want[key]) * 0.12)
                out.append({"pid": m["pid"], "check": key, "got": round(got, 4),
                            "want": want[key], "ok": ok})
        for key, got in m.get("contrasts", {}).items():
            floor = m.get("min_contrast", {}).get(key, 4.5)
            out.append({"pid": m["pid"], "check": f"contraste:{key}", "got": round(got, 2),
                        "want": f">= {floor}", "ok": got >= floor})
        for key, pt in m.get("type_pt", {}).items():
            if key in ("kicker", "label"):
                continue  # micro-capitales : elles ont leur propre plancher, cf. scale()
            out.append({"pid": m["pid"], "check": f"plancher:{key}", "got": pt,
                        "want": f">= {FLOOR_BODY_PT}", "ok": pt >= FLOOR_BODY_PT})
        # Un chanfrein posé sur une carte trop haute redevient invisible : la coupe doit peser
        # au moins 5 % de la PETITE dimension, sinon elle ne se lit plus comme une signature.
        if "notch_over_width" in m["measured"]:
            short = min(m["w"], m["h"])
            got = m["measured"]["notch_over_width"] * m["w"] / short
            out.append({"pid": m["pid"], "check": "chanfrein/petite dimension", "got": round(got, 3),
                        "want": "0.05 – 0.13", "ok": 0.05 <= got <= 0.13})

    bad = [o for o in out if not o["ok"]]
    if bad and raise_on_fail:
        lines = "\n".join(
            f"  ✗ {o['pid']} · {o['check']} : mesuré {o['got']}, attendu {o['want']}" for o in bad
        )
        raise AssertionError(f"{len(bad)} écart(s) aux benchmarks de la bibliothèque :\n{lines}")
    return out
